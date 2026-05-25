"""
presentations/pptx_builder.py
INSTAASYS PPTX Builder — Reference-quality slide generation.
Supports: title, objectives, bullets, cards, stats, section, quote, comparison, summary.
"""

import io
import logging
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════
# CONTENT REPAIR — normalises paragraph-shaped AI output into
# proper bullet arrays before any slide is rendered.
# ═══════════════════════════════════════════════════════════

def _repair_content(content, layout='bullets'):
    """Normalise AI content into a proper list of short bullet strings."""
    if not content:
        return []

    # Object-shaped layouts pass through unchanged
    if layout in ('cards', 'stats', 'comparison', 'quote'):
        return content
    if content and isinstance(content[0], dict):
        return content

    # Flatten nested lists
    flat = []
    for item in content:
        if isinstance(item, list):
            flat.extend(item)
        else:
            flat.append(str(item))

    # Single long paragraph → split it
    if len(flat) == 1 and len(flat[0]) > 80:
        paragraph = flat[0]

        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', paragraph)
        if len(sentences) >= 2:
            bullets = []
            for s in sentences:
                s = s.strip().rstrip('.')
                if len(s) > 15:
                    bullets.append(s)
            result = []
            for b in bullets[:6]:
                words = b.split()
                if len(words) > 30:
                    b = ' '.join(words[:28]) + '...'
                result.append(b)
            return result if result else [paragraph[:120] + '...']

        clauses = [c.strip() for c in paragraph.split(',') if len(c.strip()) > 15]
        if len(clauses) >= 2:
            bullets = []
            current = ''
            for clause in clauses:
                if not current:
                    current = clause
                elif len(current) + len(clause) < 80:
                    current += ', ' + clause
                else:
                    bullets.append(current)
                    current = clause
            if current:
                bullets.append(current)
            return bullets[:6] if bullets else flat

    # Multiple items, but some too long → truncate
    result = []
    for item in flat[:6]:
        words = str(item).split()
        if len(words) > 30:
            item = ' '.join(words[:28]) + '...'
        result.append(item)
    return result

# ═══════════════════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════════════════
COLORS = {
    'navy':       RGBColor(0x0F, 0x17, 0x2A),
    'navy_mid':   RGBColor(0x1E, 0x3A, 0x5F),
    'accent':     RGBColor(0x3B, 0x82, 0xF6),
    'white':      RGBColor(0xFF, 0xFF, 0xFF),
    'dark_text':  RGBColor(0x1E, 0x29, 0x3B),
    'muted':      RGBColor(0x64, 0x74, 0x8B),
    'light_bg':   RGBColor(0xF8, 0xFA, 0xFC),
    'border':     RGBColor(0xE2, 0xE8, 0xF0),
    'footer_bg':  RGBColor(0x1E, 0x29, 0x3B),
    'card_1':     RGBColor(0x23, 0x5A, 0xD0),
    'card_2':     RGBColor(0x0D, 0x7C, 0x6C),
    'card_3':     RGBColor(0xC0, 0x51, 0x0A),
    'card_4':     RGBColor(0x6D, 0x28, 0xD9),
    'card_5':     RGBColor(0xB9, 0x1C, 0x1C),
    'card_6':     RGBColor(0x16, 0x65, 0x34),
}

CARD_COLORS = [
    COLORS['card_1'], COLORS['card_2'],
    COLORS['card_3'], COLORS['card_4'],
    COLORS['card_5'], COLORS['card_6'],
]

FONT_FACE = 'Calibri'

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def repair_ai_data(ai_data: dict) -> dict:
    """Repair AI slide data in place: convert paragraph-shaped content to bullet
    arrays. Skip layouts whose content is object-shaped or naturally text-only."""
    if not ai_data:
        return ai_data
    skip_layouts = {'cards', 'stats', 'comparison', 'quote', 'section', 'title'}
    for slide in ai_data.get('slides', []) or []:
        layout = (slide.get('layout') or '').lower()
        if layout in skip_layouts:
            continue
        slide['content'] = _repair_content(slide.get('content', []), layout or 'bullets')
    return ai_data


# ═══════════════════════════════════════════════════════════
# PUBLIC ENTRY POINT
# ═══════════════════════════════════════════════════════════

def build_pptx(ai_data: dict, course=None, collect_stats: bool = False):
    """
    Build a .pptx from ai_data and return raw bytes.
    When collect_stats=True returns (bytes, stats_dict).
    Accepts both old and new layout field names.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = ai_data.get('slides', [])

    for i, slide_data in enumerate(slides):
        # Normalise old field names so both pipeline versions work
        slide_data = _normalise_slide(slide_data)
        layout = slide_data.get('layout', 'bullets').lower()
        slide_type = _detect_layout(slide_data, i, layout)

        if slide_type == 'title':
            _build_title_slide(prs, slide_data, course, ai_data)
        elif slide_type == 'objectives':
            _build_objectives_slide(prs, slide_data, course)
        elif slide_type == 'section':
            _build_section_slide(prs, slide_data, course)
        elif slide_type == 'cards':
            _build_cards_slide(prs, slide_data, course)
        elif slide_type == 'stats':
            _build_stats_slide(prs, slide_data, course)
        elif slide_type == 'quote':
            _build_quote_slide(prs, slide_data, course)
        elif slide_type == 'comparison':
            _build_comparison_slide(prs, slide_data, course)
        elif slide_type == 'summary':
            _build_summary_slide(prs, slide_data, course)
        else:
            _build_bullets_slide(prs, slide_data, course)

    buf = io.BytesIO()
    prs.save(buf)
    raw = buf.getvalue()

    if collect_stats:
        stats = {
            'slides_total': len(slides),
            'images_fetched': 0,
            'images_failed': 0,
            'image_sources': {},
        }
        return raw, stats
    return raw


def _normalise_slide(slide_data: dict) -> dict:
    """Normalise field name variants from different pipeline versions."""
    d = dict(slide_data)
    if not d.get('layout'):
        d['layout'] = d.pop('layout_type', '') or d.pop('visual_style', '') or 'bullets'
    # Map old layout names to new ones
    layout_map = {
        'content': 'bullets',
        'conclusion': 'summary',
        'recap': 'summary',
        'intro': 'bullets',
        'introduction': 'bullets',
        'background': 'bullets',
        'outline': 'bullets',
    }
    d['layout'] = layout_map.get(d['layout'], d['layout'])
    return d


def _detect_layout(slide_data, idx, layout):
    """Determine the actual layout type to render."""
    if idx == 0 or layout == 'title':
        return 'title'
    if layout == 'objectives' or 'objective' in slide_data.get('title', '').lower():
        return 'objectives'
    if layout == 'section':
        return 'section'
    if layout == 'cards':
        content = slide_data.get('content', [])
        if content and isinstance(content[0], dict) and 'heading' in content[0]:
            return 'cards'
    if layout == 'stats':
        content = slide_data.get('content', [])
        if content and isinstance(content[0], dict) and 'number' in content[0]:
            return 'stats'
    if layout == 'quote':
        content = slide_data.get('content', [])
        if content and isinstance(content[0], dict) and 'text' in content[0]:
            return 'quote'
    if layout == 'comparison':
        content = slide_data.get('content', [])
        if content and isinstance(content[0], dict) and 'left_title' in content[0]:
            return 'comparison'
    title = slide_data.get('title', '').lower()
    if layout == 'summary' or 'summary' in title or 'takeaway' in title or 'conclusion' in title:
        return 'summary'
    return 'bullets'


# ═══════════════════════════════════════════════════════════
# SHARED HELPERS
# ═══════════════════════════════════════════════════════════

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _course_code(course):
    return course.code if course and hasattr(course, 'code') else str(course or '')


def _strip_week_prefix(value):
    """Return week as a bare string, removing any leading 'Week ' so callers can
    safely render f'Week {value}' without producing 'Week Week 1'."""
    s = str(value or '').strip()
    if s.lower().startswith('week '):
        s = s[5:].strip()
    return s


def _add_footer(slide, course, week_number):
    footer = slide.shapes.add_shape(
        1, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = COLORS['footer_bg']
    footer.line.fill.background()

    left_tb = slide.shapes.add_textbox(Inches(0.25), SLIDE_H - Inches(0.35), Inches(3), Inches(0.35))
    p = left_tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = _course_code(course)
    r.font.size = Pt(9)
    r.font.color.rgb = COLORS['muted']
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.LEFT

    mid_tb = slide.shapes.add_textbox(Inches(5), SLIDE_H - Inches(0.35), Inches(3.33), Inches(0.35))
    p = mid_tb.text_frame.paragraphs[0]
    r = p.add_run()
    week_str = _strip_week_prefix(week_number)
    r.text = f'Week {week_str}' if week_str else ''
    r.font.size = Pt(9)
    r.font.color.rgb = COLORS['muted']
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.CENTER

    accent = slide.shapes.add_shape(1, SLIDE_W - Inches(0.6), SLIDE_H - Inches(0.32), Inches(0.5), Inches(0.22))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent']
    accent.line.fill.background()


def _add_header_bar(slide, title, course_code_str):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['navy']
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(11.2), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = COLORS['white']
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.LEFT

    badge = slide.shapes.add_textbox(Inches(11.8), Inches(0.2), Inches(1.3), Inches(0.5))
    p = badge.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = course_code_str
    r.font.size = Pt(11)
    r.font.color.rgb = COLORS['accent']
    r.font.bold = True
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.RIGHT


def _set_paragraph_bullet(para, char='•'):
    pPr = para._pPr
    if pPr is None:
        pPr = etree.SubElement(para._p, qn('a:pPr'))
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', char)


def _add_run(para, text, size=None, bold=False, color=None, font=None):
    r = para.add_run()
    r.text = text
    r.font.name = font or FONT_FACE
    if size:
        r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return r


def _add_speaker_notes(slide, notes_text: str):
    if not notes_text:
        return
    if not slide.has_notes_slide:
        slide.notes_slide  # creates it
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes_text


# ═══════════════════════════════════════════════════════════
# TITLE SLIDE
# ═══════════════════════════════════════════════════════════

def _build_title_slide(prs, data, course, ai_data=None):
    slide = _blank_slide(prs)

    left = slide.shapes.add_shape(1, 0, 0, Inches(3.8), SLIDE_H)
    left.fill.solid()
    left.fill.fore_color.rgb = COLORS['navy']
    left.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()

    code = _course_code(course)
    pill = slide.shapes.add_shape(1, Inches(0.4), Inches(0.5), Inches(1.4), Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLORS['accent']
    pill.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.51), Inches(1.38), Inches(0.42))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = COLORS['white']
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.CENTER

    title_text = (ai_data or {}).get('title') or data.get('title', '')
    if not title_text and course:
        title_text = getattr(course, 'title', '')
    tb_title = slide.shapes.add_textbox(Inches(4.1), Inches(1.5), Inches(8.8), Inches(3.5))
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = COLORS['dark_text']
    r.font.name = FONT_FACE

    week = _strip_week_prefix(data.get('week', '') or ((ai_data or {}).get('week', '')))
    semester = getattr(course, 'semester', '') if course else ''
    sy = getattr(course, 'school_year', '') if course else ''
    info_lines = []
    if week:
        info_lines.append(f'Week {week}')
    if semester:
        info_lines.append(f'{semester} Semester')
    if sy:
        info_lines.append(f'S.Y. {sy}')

    tb_info = slide.shapes.add_textbox(Inches(4.1), Inches(5.2), Inches(8.8), Inches(1.2))
    tf = tb_info.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = '  ·  '.join(info_lines)
    r.font.size = Pt(13)
    r.font.color.rgb = COLORS['muted']
    r.font.name = FONT_FACE

    foot = slide.shapes.add_shape(1, 0, SLIDE_H - Inches(0.6), SLIDE_W, Inches(0.6))
    foot.fill.solid()
    foot.fill.fore_color.rgb = COLORS['footer_bg']
    foot.line.fill.background()

    instructor = getattr(course, 'instructor', None) if course else None
    if instructor:
        name = instructor.get_full_name() or instructor.username
    else:
        name = ''

    tb_foot = slide.shapes.add_textbox(Inches(0.25), SLIDE_H - Inches(0.58), Inches(9), Inches(0.55))
    p = tb_foot.text_frame.paragraphs[0]
    _add_run(p, name, size=11, color=COLORS['white'])

    course_title = getattr(course, 'title', '') if course else ''
    short_title = (course_title[:50] + '...') if len(course_title) > 50 else course_title
    tb_ct = slide.shapes.add_textbox(Inches(9.5), SLIDE_H - Inches(0.58), Inches(3.5), Inches(0.55))
    p = tb_ct.text_frame.paragraphs[0]
    _add_run(p, short_title, size=10, color=COLORS['muted'])

    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# OBJECTIVES SLIDE
# ═══════════════════════════════════════════════════════════

def _build_objectives_slide(prs, data, course):
    data = dict(data)
    data['content'] = _repair_content(data.get('content', []), 'objectives')

    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_bg']
    bg.line.fill.background()

    _add_header_bar(slide, data.get('title', 'Learning Objectives'), _course_code(course))
    _add_footer(slide, course, data.get('week', ''))

    objectives = data.get('content', [])
    if not objectives:
        _add_speaker_notes(slide, data.get('notes', ''))
        return

    n = len(objectives)
    top_start = Inches(1.3)
    bottom_margin = Inches(0.5)
    available_h = SLIDE_H - top_start - bottom_margin
    item_h = available_h / n  # fill the slide regardless of count

    if n <= 3:
        font_size = 17
    elif n <= 4:
        font_size = 15
    elif n <= 5:
        font_size = 14
    else:
        font_size = 13

    circle_size = Inches(0.5)

    for i, obj in enumerate(objectives):
        y = top_start + i * item_h
        circle_y = y + (item_h / 2) - (circle_size / 2)

        # Numbered circle on the left
        circle = slide.shapes.add_shape(9, Inches(0.35), circle_y, circle_size, circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['accent']
        circle.line.fill.background()

        tb_num = slide.shapes.add_textbox(Inches(0.35), circle_y, circle_size, circle_size)
        p = tb_num.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = COLORS['white']
        r.font.name = FONT_FACE
        p.alignment = PP_ALIGN.CENTER

        # White card row with subtle border
        card_bg = slide.shapes.add_shape(
            1, Inches(1.0), y + Inches(0.05),
            Inches(12.0), item_h - Inches(0.1)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLORS['white']
        card_bg.line.color.rgb = COLORS['border']
        card_bg.line.width = Pt(0.5)

        # Left accent stripe on card
        accent_bar = slide.shapes.add_shape(
            1, Inches(1.0), y + Inches(0.05),
            Inches(0.06), item_h - Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLORS['accent']
        accent_bar.line.fill.background()

        tb_obj = slide.shapes.add_textbox(
            Inches(1.2), y + Inches(0.08),
            Inches(11.6), item_h - Inches(0.16)
        )
        tf = tb_obj.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = obj if isinstance(obj, str) else str(obj)
        r.font.size = Pt(font_size)
        r.font.color.rgb = COLORS['dark_text']
        r.font.name = FONT_FACE
        p.alignment = PP_ALIGN.LEFT

    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# BULLETS SLIDE
# ═══════════════════════════════════════════════════════════

def _build_bullets_slide(prs, data, course):
    data = dict(data)
    data['content'] = _repair_content(data.get('content', []), 'bullets')

    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_bg']
    bg.line.fill.background()

    _add_header_bar(slide, data.get('title', ''), _course_code(course))
    _add_footer(slide, course, data.get('week', ''))

    bullets = data.get('content', [])
    if not bullets:
        _add_speaker_notes(slide, data.get('notes', ''))
        return

    n = len(bullets)
    top_start = Inches(1.25)
    available_h = SLIDE_H - top_start - Inches(0.5)
    row_h = available_h / n
    left_margin = Inches(0.4)
    row_w = SLIDE_W - left_margin - Inches(0.4)

    font_size = 16 if n <= 3 else 14 if n <= 4 else 13 if n <= 5 else 12

    for i, bullet in enumerate(bullets):
        y = top_start + i * row_h
        bullet_str = bullet if isinstance(bullet, str) else str(bullet)

        # Subtle white row card on every other bullet for rhythm
        if i % 2 == 0:
            row_bg = slide.shapes.add_shape(
                1, left_margin, y + Inches(0.04),
                row_w, row_h - Inches(0.08)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = COLORS['white']
            row_bg.line.color.rgb = COLORS['border']
            row_bg.line.width = Pt(0.3)

        # Accent dot
        dot = slide.shapes.add_shape(
            9, left_margin + Inches(0.1),
            y + (row_h / 2) - Inches(0.1),
            Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLORS['accent']
        dot.line.fill.background()

        tb = slide.shapes.add_textbox(
            left_margin + Inches(0.45),
            y + Inches(0.06),
            row_w - Inches(0.55),
            row_h - Inches(0.12)
        )
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]

        if ' — ' in bullet_str:
            heading, desc = bullet_str.split(' — ', 1)
            r1 = p.add_run()
            r1.text = heading + ' — '
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = COLORS['navy']
            r1.font.name = FONT_FACE

            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(font_size - 1)
            r2.font.color.rgb = COLORS['dark_text']
            r2.font.name = FONT_FACE
        else:
            r = p.add_run()
            r.text = bullet_str
            r.font.size = Pt(font_size)
            r.font.color.rgb = COLORS['dark_text']
            r.font.name = FONT_FACE

    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# CARDS SLIDE — 2x2 or 2x3 grid
# ═══════════════════════════════════════════════════════════

def _build_cards_slide(prs, data, course):
    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_bg']
    bg.line.fill.background()

    _add_header_bar(slide, data.get('title', ''), _course_code(course))
    _add_footer(slide, course, data.get('week', ''))

    cards = data.get('content', [])
    if not cards:
        _add_speaker_notes(slide, data.get('notes', ''))
        return

    n = len(cards)
    cols = 2
    rows = (n + 1) // 2

    margin_x = Inches(0.35)
    margin_y = Inches(1.2)
    gap_x = Inches(0.2)
    gap_y = Inches(0.18)
    available_w = SLIDE_W - 2 * margin_x - gap_x
    available_h = SLIDE_H - margin_y - Inches(0.5)
    card_w = available_w / cols
    card_h = available_h / rows - gap_y
    header_h = Inches(0.48)

    for idx, card in enumerate(cards):
        row = idx // cols
        col = idx % cols
        x = margin_x + col * (card_w + gap_x)
        y = margin_y + row * (card_h + gap_y)

        heading = card.get('heading', '') if isinstance(card, dict) else str(card)
        body = card.get('body', '') if isinstance(card, dict) else ''

        card_bg = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLORS['white']
        card_bg.line.color.rgb = COLORS['border']
        card_bg.line.width = Pt(0.5)

        color = CARD_COLORS[idx % len(CARD_COLORS)]
        header = slide.shapes.add_shape(1, x, y, card_w, header_h)
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()

        tb_h = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.05),
            card_w - Inches(0.3), header_h - Inches(0.08)
        )
        p = tb_h.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = heading
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = COLORS['white']
        r.font.name = FONT_FACE
        p.alignment = PP_ALIGN.LEFT

        if body:
            tb_b = slide.shapes.add_textbox(
                x + Inches(0.15), y + header_h + Inches(0.1),
                card_w - Inches(0.3), card_h - header_h - Inches(0.15)
            )
            tf = tb_b.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = body
            r.font.size = Pt(12 if rows <= 2 else 11)
            r.font.color.rgb = COLORS['dark_text']
            r.font.name = FONT_FACE
            p.alignment = PP_ALIGN.LEFT

    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# STATS SLIDE — large number display
# ═══════════════════════════════════════════════════════════

def _build_stats_slide(prs, data, course):
    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['navy']
    bg.line.fill.background()

    _add_footer(slide, course, data.get('week', ''))

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    p = tb_title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = data.get('title', '')
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = COLORS['white']
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.CENTER

    stats = data.get('content', [])
    if not stats:
        _add_speaker_notes(slide, data.get('notes', ''))
        return

    n = len(stats)
    stat_w = SLIDE_W / n
    y_number = Inches(1.6)
    y_label = Inches(3.4)
    y_source = Inches(4.6)

    for i, stat in enumerate(stats):
        number = stat.get('number', '') if isinstance(stat, dict) else ''
        label = stat.get('label', '') if isinstance(stat, dict) else ''
        source = stat.get('source', '') if isinstance(stat, dict) else ''
        x = i * stat_w

        tb_n = slide.shapes.add_textbox(x + Inches(0.2), y_number, stat_w - Inches(0.4), Inches(1.6))
        p = tb_n.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = number
        r.font.size = Pt(60)
        r.font.bold = True
        r.font.color.rgb = COLORS['accent']
        r.font.name = FONT_FACE
        p.alignment = PP_ALIGN.CENTER

        tb_l = slide.shapes.add_textbox(x + Inches(0.2), y_label, stat_w - Inches(0.4), Inches(1.0))
        p = tb_l.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.size = Pt(15)
        r.font.color.rgb = COLORS['white']
        r.font.name = FONT_FACE
        p.alignment = PP_ALIGN.CENTER

        if source:
            tb_s = slide.shapes.add_textbox(x + Inches(0.2), y_source, stat_w - Inches(0.4), Inches(0.6))
            p = tb_s.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = source
            r.font.size = Pt(11)
            r.font.color.rgb = COLORS['muted']
            r.font.name = FONT_FACE
            p.alignment = PP_ALIGN.CENTER

        if i < n - 1:
            div = slide.shapes.add_shape(
                1, x + stat_w - Inches(0.02), Inches(1.5),
                Inches(0.02), Inches(3.5)
            )
            div.fill.solid()
            div.fill.fore_color.rgb = RGBColor(0x33, 0x4E, 0x6B)
            div.line.fill.background()

    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# SECTION DIVIDER SLIDE
# ═══════════════════════════════════════════════════════════

def _build_section_slide(prs, data, course):
    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['navy']
    bg.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()

    section_label = data.get('section_label', 'SECTION')
    tb_label = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(0.5))
    p = tb_label.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = section_label.upper()
    r.font.size = Pt(13)
    r.font.color.rgb = COLORS['accent']
    r.font.bold = True
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.LEFT

    tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12), Inches(2.0))
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data.get('title', '')
    r.font.size = Pt(38)
    r.font.bold = True
    r.font.color.rgb = COLORS['white']
    r.font.name = FONT_FACE
    p.alignment = PP_ALIGN.LEFT

    content = data.get('content', [])
    subtitle = ''
    if content:
        first = content[0]
        if isinstance(first, str):
            subtitle = first
        elif isinstance(first, dict):
            subtitle = first.get('text', '')
    if subtitle:
        tb_sub = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12), Inches(0.8))
        p = tb_sub.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(16)
        r.font.color.rgb = COLORS['muted']
        r.font.name = FONT_FACE
        p.alignment = PP_ALIGN.LEFT

    _add_footer(slide, course, data.get('week', ''))
    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# QUOTE SLIDE
# ═══════════════════════════════════════════════════════════

def _build_quote_slide(prs, data, course):
    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['navy']
    bg.line.fill.background()

    content = data.get('content', [{}])
    item = content[0] if content else {}
    quote_text = item.get('text', '') if isinstance(item, dict) else str(item)
    attribution = item.get('attribution', '') if isinstance(item, dict) else ''

    tb_q = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(2), Inches(1.5))
    p = tb_q.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = '“'
    r.font.size = Pt(90)
    r.font.color.rgb = COLORS['accent']
    r.font.name = FONT_FACE

    tb_quote = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(3.5))
    tf = tb_quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = quote_text
    r.font.size = Pt(22)
    r.font.color.rgb = COLORS['white']
    r.font.name = FONT_FACE
    r.font.italic = True

    if attribution:
        tb_attr = slide.shapes.add_textbox(Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.6))
        p = tb_attr.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f'— {attribution}'
        r.font.size = Pt(14)
        r.font.color.rgb = COLORS['muted']
        r.font.name = FONT_FACE

    _add_header_bar(slide, data.get('title', ''), _course_code(course))
    _add_footer(slide, course, data.get('week', ''))
    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# COMPARISON SLIDE — side-by-side
# ═══════════════════════════════════════════════════════════

def _build_comparison_slide(prs, data, course):
    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_bg']
    bg.line.fill.background()

    _add_header_bar(slide, data.get('title', ''), _course_code(course))
    _add_footer(slide, course, data.get('week', ''))

    content = data.get('content', [{}])
    comp = content[0] if content else {}

    left_title = comp.get('left_title', 'A') if isinstance(comp, dict) else 'A'
    right_title = comp.get('right_title', 'B') if isinstance(comp, dict) else 'B'
    left_points = comp.get('left_points', []) if isinstance(comp, dict) else []
    right_points = comp.get('right_points', []) if isinstance(comp, dict) else []

    col_w = Inches(6.2)
    col_h = SLIDE_H - Inches(1.7)
    header_h = Inches(0.55)
    margin = Inches(0.35)
    gap = Inches(0.2)

    for side_idx, (title, points, color) in enumerate([
        (left_title, left_points, COLORS['card_1']),
        (right_title, right_points, COLORS['card_2'])
    ]):
        x = margin if side_idx == 0 else margin + col_w + gap

        card = slide.shapes.add_shape(1, x, Inches(1.2), col_w, col_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = COLORS['border']
        card.line.width = Pt(0.5)

        hdr = slide.shapes.add_shape(1, x, Inches(1.2), col_w, header_h)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()

        tb_hdr = slide.shapes.add_textbox(x + Inches(0.2), Inches(1.25), col_w - Inches(0.4), header_h)
        p = tb_hdr.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = COLORS['white']
        r.font.name = FONT_FACE

        tb_pts = slide.shapes.add_textbox(
            x + Inches(0.2), Inches(1.2) + header_h + Inches(0.1),
            col_w - Inches(0.4), col_h - header_h - Inches(0.2)
        )
        tf = tb_pts.text_frame
        tf.word_wrap = True

        for i, pt in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(7)
            _set_paragraph_bullet(p)
            r = p.add_run()
            r.text = str(pt)
            r.font.size = Pt(13)
            r.font.color.rgb = COLORS['dark_text']
            r.font.name = FONT_FACE

    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# SUMMARY / KEY TAKEAWAYS SLIDE
# ═══════════════════════════════════════════════════════════

def _build_summary_slide(prs, data, course):
    data = dict(data)
    data['content'] = _repair_content(data.get('content', []), 'summary')

    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['navy']
    bg.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()

    tb_title = slide.shapes.add_textbox(Inches(0.45), Inches(0.3), Inches(12), Inches(0.9))
    p = tb_title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = data.get('title', 'Key Takeaways')
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = COLORS['white']
    r.font.name = FONT_FACE

    underline = slide.shapes.add_shape(1, Inches(0.45), Inches(1.1), Inches(2.5), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = COLORS['accent']
    underline.line.fill.background()

    takeaways = data.get('content', [])
    n = len(takeaways)
    if n == 0:
        _add_footer(slide, course, data.get('week', ''))
        _add_speaker_notes(slide, data.get('notes', ''))
        return

    quote = data.get('quote') if isinstance(data.get('quote'), dict) else None
    top = Inches(1.3)
    bottom_reserved = Inches(1.7) if quote else Inches(1.0)
    available_h = SLIDE_H - top - bottom_reserved
    item_h = available_h / n  # evenly distributed, no fixed cap

    if n <= 4:
        font_size = 17
    elif n <= 5:
        font_size = 15
    else:
        font_size = 13

    circle_size = Inches(0.42)

    for i, point in enumerate(takeaways):
        y = top + i * item_h
        circle_y = y + (item_h / 2) - (circle_size / 2)

        circ = slide.shapes.add_shape(9, Inches(0.4), circle_y, circle_size, circle_size)
        circ.fill.solid()
        circ.fill.fore_color.rgb = COLORS['accent']
        circ.line.fill.background()

        tb_check = slide.shapes.add_textbox(Inches(0.4), circle_y, circle_size, circle_size)
        p = tb_check.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = '✓'
        r.font.size = Pt(13)
        r.font.color.rgb = COLORS['white']
        r.font.bold = True
        r.font.name = FONT_FACE
        p.alignment = PP_ALIGN.CENTER

        tb_pt = slide.shapes.add_textbox(Inches(1.05), y + Inches(0.05), Inches(12.0), item_h - Inches(0.1))
        tf = tb_pt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(point)
        r.font.size = Pt(font_size)
        r.font.color.rgb = COLORS['white']
        r.font.name = FONT_FACE

    quote = data.get('quote')
    if quote and isinstance(quote, dict):
        tb_q = slide.shapes.add_textbox(Inches(0.45), SLIDE_H - Inches(1.1), Inches(12.5), Inches(0.7))
        p = tb_q.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f'“{quote["text"]}”'
        if quote.get('attribution'):
            r.text += f' — {quote["attribution"]}'
        r.font.size = Pt(13)
        r.font.color.rgb = COLORS['muted']
        r.font.italic = True
        r.font.name = FONT_FACE

    _add_footer(slide, course, data.get('week', ''))
    _add_speaker_notes(slide, data.get('notes', ''))


# ═══════════════════════════════════════════════════════════
# BACKWARD-COMPAT: web image URL fetch used by ai_pipeline.py
# ═══════════════════════════════════════════════════════════

def fetch_image_url_for_web(query: str):
    """Return a CDN image URL for web display. Pexels → Unsplash."""
    if not query:
        return None
    try:
        from django.conf import settings as django_settings
        pexels_key = getattr(django_settings, 'PEXELS_API_KEY', '')
        unsplash_key = getattr(django_settings, 'UNSPLASH_ACCESS_KEY', '')
    except Exception:
        pexels_key = unsplash_key = ''

    try:
        import requests
        if pexels_key:
            r = requests.get(
                'https://api.pexels.com/v1/search',
                params={'query': query, 'per_page': 3, 'orientation': 'landscape'},
                headers={'Authorization': pexels_key},
                timeout=5,
            )
            photos = r.json().get('photos', [])
            if photos:
                src = photos[0].get('src', {})
                url = src.get('large') or src.get('original')
                if url:
                    return url

        if unsplash_key:
            r = requests.get(
                'https://api.unsplash.com/search/photos',
                params={'query': query, 'per_page': 3, 'orientation': 'landscape'},
                headers={'Authorization': f'Client-ID {unsplash_key}'},
                timeout=5,
            )
            results = r.json().get('results', [])
            if results:
                url = results[0].get('urls', {}).get('regular')
                if url:
                    return url
    except Exception as exc:
        logger.debug(f"Web image URL fetch failed for '{query}': {exc}")
    return None
